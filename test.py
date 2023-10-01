##test
import openai
from pptx import Presentation

openai.api_key = ""

prs = Presentation()

# Generate text using OpenAI API
response = openai.Completion.create(engine="davinci", prompt="Title: Benefits of Artificial Intelligence", max_tokens=150)
generated_text = response['choices'][0]['text'].strip()

# Split generated text into title and content
title, content = generated_text.split('\n', 1)

# Add a slide with a title and content
slide = prs.slides.add_slide(prs.slide_layouts[1])
title_shape = slide.shapes.title
title_shape.text = title
content_box = slide.placeholders[1]
content_box.text = content

# Save the presentation
prs.save('presentation.pptx')